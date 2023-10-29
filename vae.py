import torch
import torch.nn as nn
import torch.optim as optim
from sklearn.preprocessing import StandardScaler
from sklearn.model_selection import train_test_split
import pandas as pd
from torch.utils.data import DataLoader, TensorDataset
import matplotlib.pyplot as plt
import numpy as np
from scipy.stats import pearsonr
from sklearn.manifold import TSNE
from pptx import Presentation
from pptx.util import Inches



class VAE(nn.Module):
    def __init__(self, input_dim, latent_dim):
        super(VAE, self).__init__()

        # dimensions for the intermediate layers
        dims = [input_dim, int(input_dim*1.3), int(input_dim*1.2), int(input_dim*1), int(input_dim*0.5), int(input_dim*0.3), latent_dim, latent_dim]

        # Encoder layers with BatchNorm
        encoder_layers = []
        for i in range(len(dims)-1):
            encoder_layers.append(nn.Linear(dims[i], dims[i+1]))
            encoder_layers.append(nn.BatchNorm1d(dims[i+1]))
            encoder_layers.append(nn.ReLU())
        self.encoder = nn.Sequential(*encoder_layers)

        # Latent space layers
        self.fc_mu = nn.Linear(dims[-1], latent_dim)
        self.fc_logvar = nn.Linear(dims[-1], latent_dim)

        # Decoder layers with BatchNorm
        decoder_layers = []
        for i in range(len(dims)-2, -1, -1):
            decoder_layers.append(nn.Linear(dims[i+1], dims[i]))
            decoder_layers.append(nn.BatchNorm1d(dims[i]))
            decoder_layers.append(nn.ReLU())
            #print(len(decoder_layers))
        decoder_layers.append(nn.Sigmoid()) 
        self.decoder = nn.Sequential(*decoder_layers)

    def encode(self, x):
        hidden = self.encoder(x)
        mu = self.fc_mu(hidden)
        logvar = self.fc_logvar(hidden)
        return mu, logvar

    def reparameterize(self, mu, logvar):
        std = torch.exp(0.5 * logvar)
        eps = torch.randn_like(std)
        z = mu + eps * std
        return z

    def decode(self, z):
        output = self.decoder(z)
        return output

    def forward(self, x):
        mu, logvar = self.encode(x)
        z = self.reparameterize(mu, logvar)
        output = self.decode(z)
        return output, mu, logvar




def create_dataloader(train_tensor, batch_size):
    dataset = TensorDataset(train_tensor)
    dataloader = DataLoader(dataset, batch_size=batch_size, shuffle=False)
    return dataloader




def kl_divergence(mu, logvar):
    # -0.5 * sum(1 + log(sigma^2) - mu^2 - sigma^2)
    return -0.5 * torch.sum(1 + logvar - mu.pow(2) - logvar.exp())


def pearson_correlation(x, y):
    vx = x - torch.mean(x)
    vy = y - torch.mean(y)
    return torch.sum(vx * vy) / (torch.sqrt(torch.sum(vx ** 2)) * torch.sqrt(torch.sum(vy ** 2)))






filename = '/mnt/data/macaulay/datas/processed_OmicsExpression.csv'
df = pd.read_csv(filename)
cell_line_names = df['Cell_line'].copy()  # Save the 'Cell_line' column
df = df.drop(columns=['Cell_line'])

numeric_data = df.values
scaler = StandardScaler()
scaled_data = scaler.fit_transform(numeric_data)
train_tensor = torch.tensor(scaled_data, dtype=torch.float32)
print("df shape:", df.shape)
print('training tensor data shape:', train_tensor.shape)

device = ("cuda" if torch.cuda.is_available() else "cpu")
print('Using: ', device)

train_tensor = train_tensor.to(device)


input_dim = train_tensor.shape[1]
learning_rate = 0.00001
batch_size = 16
num_epochs = 1000

latent_dim = 1024



prs = Presentation()
vae = VAE(input_dim, latent_dim)
vae = vae.to(device)
reconstruction_loss = nn.MSELoss()
optimizer = optim.Adam(vae.parameters(), lr=learning_rate)
dataloader = create_dataloader(train_tensor, batch_size)

reconst_loss = []
train_losses = []

kl_losses = []
correlations = []
for epoch in range(num_epochs):
    embeddings_list = []
    vae.train()
    running_recon_loss = 0.0
    running_kl_loss = 0.0
    reconstructed_data = []

    for batch_data in dataloader:
        optimizer.zero_grad()
        batch_data = batch_data[0].to(device)

        # Forward pass
        output, mu, logvar = vae(batch_data)
        reconstructed_data.append(output.detach())

        recon_loss = reconstruction_loss(output, batch_data)
        kl_loss = kl_divergence(mu, logvar)
 
        # Combined loss
        combined_loss = recon_loss + kl_loss

        # Backward pass and optimization
        combined_loss.backward()
        optimizer.step()

        # Accumulate running losses
        running_recon_loss += recon_loss.item()
        running_kl_loss += kl_loss.item()
        embeddings = mu.detach().cpu().numpy()
        embeddings_list.append(embeddings)


    # After all batches are processed, compute correlation over the entire dataset
    reconstructed_data = torch.cat(reconstructed_data, dim=0)
    epoch_correlation = pearson_correlation(reconstructed_data, train_tensor).item()
    correlations.append(epoch_correlation)

    # Print epoch statistics
    epoch_recon_loss = running_recon_loss / len(dataloader)
    epoch_kl_loss = running_kl_loss / len(dataloader)
    epoch_combined_loss = epoch_recon_loss + epoch_kl_loss

    print(f"Epoch [{epoch + 1}/{num_epochs}], Combined Training Loss: {epoch_combined_loss:.4f}, Reconstruction Loss: {epoch_recon_loss:.4f}, KL Loss: {epoch_kl_loss:.4f}, Correlation: {epoch_correlation:.4f}")

    reconst_loss.append(epoch_recon_loss)
    train_losses.append(epoch_combined_loss)
    kl_losses.append(epoch_kl_loss)


    if epoch % 10 == 0:


        vae.eval()
        

        # Encode the entire dataset to get the latent variables
        with torch.no_grad():
            mu, _ = vae.encode(train_tensor)

        mu = mu.cpu().numpy()

        # Apply t-SNE
        tsne = TSNE(n_components=2, random_state=42)
        latent_tsne = tsne.fit_transform(mu)

        # Plot t-SNE results
        plt.figure(figsize=(5, 5))
        plt.scatter(latent_tsne[:, 0], latent_tsne[:, 1], alpha=0.5)
        plt.xlabel('t-SNE 1')
        plt.ylabel('t-SNE 2')
        plt.title(f'Latent Space t-SNE Visualization of epoch {epoch+1}')
        tmp_img_path = f'/mnt/data/macaulay/plot_images/tnse{epoch}.png'
        plt.savefig(tmp_img_path)
        plt.close()

        # Add a slide and insert the image
        slide_layout = prs.slide_layouts[5]  # Blank slide layout
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.add_picture(tmp_img_path, Inches(1), Inches(1), width=Inches(5.5))

    with open('/mnt/data/macaulay/datas/vae_expression.txt', 'a') as f:
        if epoch == 0:
            f.write(f'Epoch,Combined Training Loss,Reconstruction Loss,KL Loss,Correlation\n')
            f.write(f'{epoch+1},{epoch_combined_loss:.4f},{epoch_recon_loss:.4f},{epoch_kl_loss:.4f},{epoch_correlation:.4f}\n')
        else:
            f.write(f'{epoch+1},{epoch_combined_loss:.4f},{epoch_recon_loss:.4f},{epoch_kl_loss:.4f},{epoch_correlation:.4f}\n')
    if epoch % 10 == 0:
        embeddings_array = np.concatenate(embeddings_list)
        df_embeddings = pd.DataFrame(embeddings_array)

        df_embeddings = pd.concat([cell_line_names, df_embeddings], axis=1) #join the cell line names to the embeddings
        df_embeddings.to_csv(f'/mnt/data/macaulay/datas/OmicExpression_embeddings/OmicExpression_embeddings_{epoch}.csv', index=False)


# Save the updated VAE model
torch.save(vae.state_dict(), '/mnt/data/macaulay/vae_state/vae_expression.pth')

plt.plot(range(2, num_epochs + 1), reconst_loss[1:], label='Combined Training Loss')
plt.plot(range(2, num_epochs + 1), kl_losses[1:], label='KL Divergence')
plt.xlabel('Epoch')
plt.ylabel('Loss')
plt.title('Training Loss Curves')
plt.legend()
plt.savefig('/mnt/data/macaulay/plot_images/vae_loss_curve.png')


prs.save("/mnt/data/macaulay/datas/latent_space_visualization.pptx")
